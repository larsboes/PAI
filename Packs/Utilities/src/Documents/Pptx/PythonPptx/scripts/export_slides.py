"""Export PPTX slides to PNG images via LibreOffice headless + pdftoppm.

Usage (standalone):
    python export_slides.py deck.pptx [output_dir] [--slides 0,2,4] [--dpi 150]

Usage (from pptx_helpers):
    sb = SlideBuilder("deck.pptx")
    images = sb.export_slides()                    # All slides
    images = sb.export_slides(slides=[0, 5, 17])   # Specific slides
"""

import os
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path


def export_slides(
    pptx_path: str,
    output_dir: str = "/tmp/pptx-preview",
    slides: list[int] | None = None,
    dpi: int = 150,
) -> list[str]:
    """Export PPTX slides to PNG.

    Args:
        pptx_path: Path to .pptx file.
        output_dir: Output directory for PNGs.
        slides: 0-based PPTX slide indices to export (None = all visible).
        dpi: Resolution (150 = ~2000x1125 for widescreen).

    Returns:
        Sorted list of paths to exported PNG files.

    Note:
        LibreOffice skips hidden slides. This function builds a mapping
        from PPTX slide index → PDF page number to handle this correctly.
    """
    from pptx import Presentation as _Prs

    pptx_path = str(Path(pptx_path).resolve())
    output_dir = str(Path(output_dir).resolve())
    os.makedirs(output_dir, exist_ok=True)

    # Clean old exports
    for f in Path(output_dir).glob("slide-*.png"):
        f.unlink()

    # Build slide index → PDF page mapping (LO skips hidden slides)
    prs = _Prs(pptx_path)
    idx_to_page = {}
    page = 1
    for i, slide in enumerate(prs.slides):
        if slide._element.get('show') != '0':
            idx_to_page[i] = page
            page += 1

    with tempfile.TemporaryDirectory() as tmp:
        # PPTX → PDF
        result = subprocess.run(
            ["libreoffice", "--headless", "--convert-to", "pdf",
             "--outdir", tmp, pptx_path],
            capture_output=True, text=True, timeout=120,
        )
        if result.returncode != 0:
            raise RuntimeError(f"LibreOffice conversion failed: {result.stderr}")

        pdf_path = next(Path(tmp).glob("*.pdf"))

        # PDF → PNG
        subprocess.run(
            ["pdftoppm", "-png", "-r", str(dpi), str(pdf_path),
             str(Path(output_dir) / "slide")],
            check=True, capture_output=True, timeout=120,
        )

    # Collect exported files
    all_images = sorted(Path(output_dir).glob("slide-*.png"))

    if slides is not None:
        selected = []
        for idx in slides:
            if idx not in idx_to_page:
                continue  # Hidden slide, skip
            page_num = idx_to_page[idx]
            candidates = [p for p in all_images
                          if int(p.stem.split("-")[-1]) == page_num]
            selected.extend(candidates)
        return [str(p) for p in selected]

    return [str(p) for p in all_images]


if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(description="Export PPTX slides to PNG")
    parser.add_argument("pptx", help="Path to .pptx file")
    parser.add_argument("output", nargs="?", default="/tmp/pptx-preview")
    parser.add_argument("--slides", help="0-based indices, e.g. 0,2,4")
    parser.add_argument("--dpi", type=int, default=150)
    args = parser.parse_args()

    slide_list = [int(s) for s in args.slides.split(",")] if args.slides else None
    images = export_slides(args.pptx, args.output, slide_list, args.dpi)
    print(f"Exported {len(images)} slides to {args.output}/")
    for img in images:
        print(f"  {img}")
