"""pptx_helpers.py — Reusable python-pptx primitives and recipes for slide generation.

Provides DT Telekom-themed shape helpers, layout utilities, slide recipes, and
slide management. Analyzed from the "Telekom Liquid Master" theme deck.

Classes:
    SlideBuilder  — Shape primitives, tables, images, bullets, speaker notes, find/replace
    SlideRecipes  — High-level slide templates (title, KPI, comparison, agenda, etc.)
    GanttBuilder  — Project timeline / Gantt chart slides
    TELEKOM       — Color palette constants

Usage:
    from pptx_helpers import SlideBuilder, SlideRecipes, GanttBuilder, TELEKOM

    sb = SlideBuilder("path/to/deck.pptx")
    recipes = SlideRecipes(sb)

    recipes.title_slide("My Presentation", "Lars Boes — Feb 2026")
    recipes.kpi_dashboard([("Users", "1.2M", "+12%"), ("Revenue", "€4.2M", "+8%")])
    recipes.comparison_slide("Options", "Option A", ["Fast", "Cheap"], "Option B", ["Reliable", "Scalable"])

    # Or use primitives directly:
    slide = sb.add_blank_slide()
    sb.add_table(slide, 0.5, 1.5, 12, 3, [["Name", "Score"], ["Alice", "95"]])
    sb.add_image(slide, "chart.png", 8, 1, w=4)
    sb.add_bullet_list(slide, 0.5, 1, 5, 3, ["Point A", "Point B", "Point C"])
    sb.set_speaker_notes(slide, "Talk about these metrics...")

    sb.save()
"""

from pathlib import Path
from copy import deepcopy

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree


# ============================================================
# COLOR PALETTES
# ============================================================

class TelekomColors:
    """Official DT Telekom CI colors — extracted from Telekom Liquid Master theme."""
    # Brand
    MAGENTA       = RGBColor(0xE2, 0x00, 0x74)
    MAGENTA_DARK  = RGBColor(0x9A, 0x00, 0x50)
    MAGENTA_LIGHT = RGBColor(0xF9, 0xA8, 0xD4)
    MAGENTA_10    = RGBColor(0xFC, 0xE7, 0xF3)
    # Theme colors (from Telekom Liquid Master)
    DARK          = RGBColor(0x26, 0x26, 0x26)  # near-black text
    GRAY          = RGBColor(0x6B, 0x72, 0x80)  # secondary text
    LIGHT_GRAY    = RGBColor(0xE5, 0xE5, 0xE5)  # borders, dividers
    SURFACE       = RGBColor(0xF5, 0xF5, 0xF7)  # card/column bg
    CARD_BG       = RGBColor(0xED, 0xED, 0xED)  # used on slide 10/11/13 cards
    WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
    # Theme accents (from clrScheme)
    TEAL          = RGBColor(0x32, 0xB9, 0xAF)  # accent1
    LIGHT_BLUE    = RGBColor(0xA4, 0xDE, 0xEE)  # accent2
    PEACH         = RGBColor(0xEC, 0xCC, 0xBF)  # accent3
    KHAKI         = RGBColor(0xF0, 0xE6, 0x8C)  # accent4
    SKY_BLUE      = RGBColor(0x00, 0xA8, 0xE6)  # accent5
    PURPLE        = RGBColor(0x6E, 0x64, 0x8C)  # accent6
    # Semantic
    GREEN         = RGBColor(0x05, 0x96, 0x69)
    RED           = RGBColor(0xDC, 0x26, 0x26)
    BLUE          = RGBColor(0x25, 0x63, 0xEB)
    WARNING       = RGBColor(0xD9, 0x77, 0x06)
    ERROR_RED     = RGBColor(0xED, 0x00, 0x00)  # used in slide 15 "CONTEXT FULL"
    SUCCESS_GREEN = RGBColor(0x74, 0xE2, 0x81)  # used in slide 16 agent result
    HIGHLIGHT_PINK= RGBColor(0xFC, 0xE3, 0xEB)  # used for output highlights
    NOTE_RED      = RGBColor(0xA9, 0x49, 0x49)  # used for "Downsides" cards

TELEKOM = TelekomColors()

# Font names from theme
FONT_HEADING = 'TeleNeo Office'        # majorFont: TeleNeo Office ExtraBold
FONT_HEADING_BOLD = 'TeleNeo Office'   # Use with bold=True for ExtraBold weight
FONT_BODY = 'TeleNeo Office'           # minorFont
FONT_FALLBACK = 'Arial'                # When TeleNeo not installed


# ============================================================
# XML HELPERS
# ============================================================

def _make_element(tag, attribs=None, nsmap=None):
    """Create an lxml element with proper namespace."""
    return etree.SubElement(etree.Element('dummy'), qn(tag), attrib=attribs or {})


def _set_corner_radius(shape, radius=16667):
    """Set rounded rectangle corner radius via XML.
    Default 16667 matches the deck's existing rounded rects.
    Range: 0 (square) to 50000 (full pill).
    """
    prstGeom = shape._element.spPr.find(qn('a:prstGeom'))
    if prstGeom is not None:
        avLst = prstGeom.find(qn('a:avLst'))
        if avLst is None:
            avLst = etree.SubElement(prstGeom, qn('a:avLst'))
        else:
            for child in list(avLst):
                avLst.remove(child)
        etree.SubElement(avLst, qn('a:gd'), attrib={
            'name': 'adj', 'fmla': f'val {radius}'
        })


def _apply_magenta_gradient(shape):
    """Apply the signature DT Magenta gradient fill.

    Uses scheme color tx2 (Magenta E20074) with lumMod/lumOff stops,
    matching the deck's native gradient cards (slide 6).
    Inserts gradFill BEFORE a:ln (OOXML element order matters).
    """
    sp = shape._element
    spPr = sp.spPr

    # Remove any existing fill from spPr (but keep p:style intact)
    for fill_tag in ['a:solidFill', 'a:noFill', 'a:gradFill', 'a:pattFill']:
        existing = spPr.find(qn(fill_tag))
        if existing is not None:
            spPr.remove(existing)

    # Build gradient element
    gradFill = etree.Element(qn('a:gradFill'))
    gsLst = etree.SubElement(gradFill, qn('a:gsLst'))

    # 4 stops using scheme color tx2 (matches slide 6 deck DNA)
    stops = [
        (0,     '20000', '80000'),  # lightest pink
        (34000, '40000', '60000'),
        (64000, '60000', '40000'),
        (89000, '60000', '40000'),  # medium magenta
    ]

    for pos, lum_mod, lum_off in stops:
        gs = etree.SubElement(gsLst, qn('a:gs'), attrib={'pos': str(pos)})
        clr = etree.SubElement(gs, qn('a:schemeClr'), attrib={'val': 'tx2'})
        etree.SubElement(clr, qn('a:lumMod'), attrib={'val': lum_mod})
        etree.SubElement(clr, qn('a:lumOff'), attrib={'val': lum_off})

    etree.SubElement(gradFill, qn('a:lin'), attrib={
        'ang': '5400000', 'scaled': '1'
    })

    # Insert gradFill BEFORE a:ln (OOXML requires fill before line)
    ln = spPr.find(qn('a:ln'))
    if ln is not None:
        spPr.insert(list(spPr).index(ln), gradFill)
    else:
        spPr.append(gradFill)


def _apply_shadow(shape, blur_emu=50800, dist_emu=38100, direction=5400000,
                  color='000000', opacity=40):
    """Apply an outer shadow to a shape."""
    spPr = shape._element.spPr
    effectLst = spPr.find(qn('a:effectLst'))
    if effectLst is None:
        effectLst = etree.SubElement(spPr, qn('a:effectLst'))

    outerShdw = etree.SubElement(effectLst, qn('a:outerShdw'), attrib={
        'blurRad': str(blur_emu),
        'dist': str(dist_emu),
        'dir': str(direction),
        'algn': 'tl',
    })
    srgbClr = etree.SubElement(outerShdw, qn('a:srgbClr'), attrib={'val': color})
    etree.SubElement(srgbClr, qn('a:alpha'), attrib={'val': f'{opacity}000'})


# ============================================================
# SLIDE BUILDER
# ============================================================

class SlideBuilder:
    """High-level builder for adding/editing slides in an existing PPTX."""

    def __init__(self, path: str):
        self.path = path
        self.prs = Presentation(path)

    @property
    def slide_width(self):
        return self.prs.slide_width

    @property
    def slide_height(self):
        return self.prs.slide_height

    def save(self, path: str = None):
        """Save the presentation. Defaults to overwriting the source file."""
        self.prs.save(path or self.path)

    def export_slides(self, slides=None, output_dir="/tmp/pptx-preview", dpi=150):
        """Export slides to PNG via LibreOffice + pdftoppm. Auto-saves first.

        Args:
            slides: list of 0-based indices (None = all).
            output_dir: output directory for PNGs.
            dpi: resolution (150 = ~2000x1125 for widescreen).
        Returns:
            list of PNG file paths.
        """
        import importlib, sys as _sys
        _sys.path.insert(0, str(Path(__file__).parent))
        _mod = importlib.import_module('export_slides')
        self.save()
        return _mod.export_slides(self.path, output_dir, slides, dpi)

    # --- Slide management ---

    def add_blank_slide(self):
        """Add a truly blank slide (fewest placeholders, all removed)."""
        best_layout = min(self.prs.slide_layouts,
                          key=lambda l: len(l.placeholders))
        slide = self.prs.slides.add_slide(best_layout)
        for ph in list(slide.placeholders):
            ph._element.getparent().remove(ph._element)
        for shape in list(slide.shapes):
            if hasattr(shape, 'text') and shape.text:
                txt = shape.text.lower()
                if any(kw in txt for kw in ['trennerseite', 'bild einfügen',
                                              'click to add', 'placeholder',
                                              'kapitel', 'überschrift']):
                    shape._element.getparent().remove(shape._element)
        return slide

    def add_layout_slide(self, layout_name: str):
        """Add a slide using a specific layout by name.
        Available in this deck: 'Titel lang 01', 'Titel lang 02',
        'Kapiteltrenner 01/02/03', 'Agenda', 'Titel und Inhalt',
        'Nur Titel', 'Leer', 'Zwei Inhalte', 'Drei Inhalte',
        'Vier Inhalte', 'Inhalt mit Bild', 'Bild vollflächig', etc.
        """
        layout = self.prs.slide_layouts.get_by_name(layout_name)
        if layout is None:
            raise ValueError(f"Layout '{layout_name}' not found. "
                             f"Available: {[l.name for l in self.prs.slide_layouts]}")
        return self.prs.slides.add_slide(layout)

    def remove_last_slide(self):
        """Remove the last slide from the presentation."""
        if len(self.prs.slides) > 0:
            rId = self.prs.slides._sldIdLst[-1].get(qn('r:id'))
            self.prs.part.drop_rel(rId)
            self.prs.slides._sldIdLst.remove(self.prs.slides._sldIdLst[-1])

    def remove_slides_from(self, index: int):
        """Remove all slides from index onwards (0-based)."""
        while len(self.prs.slides) > index:
            self.remove_last_slide()

    def is_hidden(self, slide) -> bool:
        """Check if a slide is hidden."""
        return slide._element.get('show') == '0'

    def set_hidden(self, slide, hidden: bool = True):
        """Hide or unhide a slide."""
        if hidden:
            slide._element.set('show', '0')
        else:
            # Remove attribute = visible (default)
            if 'show' in slide._element.attrib:
                del slide._element.attrib['show']

    def list_slides(self, include_hidden=True):
        """List all slides with index, title, and hidden status.

        Returns: list of dicts with keys: index, title, hidden.
        """
        result = []
        for i, slide in enumerate(self.prs.slides):
            hidden = self.is_hidden(slide)
            if not include_hidden and hidden:
                continue
            # Try to extract title from shapes
            title = ''
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text.strip():
                    title = shape.text.strip()[:80]
                    break
            result.append({'index': i, 'title': title, 'hidden': hidden})
        return result

    # --- Shape primitives ---

    def add_label(self, slide, x, y, w, h, text,
                  size=10, bold=False, color=None, align=PP_ALIGN.LEFT,
                  font=FONT_FALLBACK, wrap=False):
        """Add a text label. Coordinates in inches."""
        color = color or TELEKOM.DARK
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = wrap
        tf.auto_size = None
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = font
        p.alignment = align
        return tb

    def add_rect(self, slide, x, y, w, h, color=None):
        """Add a plain rectangle."""
        color = color or TELEKOM.MAGENTA
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    def add_rounded_rect(self, slide, x, y, w, h, color=None, radius=16667):
        """Add a rounded rectangle (default radius matches deck style)."""
        color = color or TELEKOM.MAGENTA
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        _set_corner_radius(shape, radius)
        return shape

    def add_rounded_bar(self, slide, x, y, w, h, color=None, radius=30000):
        """Add a rounded bar (higher radius for pill-shape, good for Gantt bars)."""
        return self.add_rounded_rect(slide, x, y, w, h, color, radius)

    def add_diamond(self, slide, x, y, size, color=None):
        """Add a diamond milestone marker."""
        color = color or TELEKOM.MAGENTA
        d = slide.shapes.add_shape(
            MSO_SHAPE.DIAMOND, Inches(x), Inches(y), Inches(size), Inches(size))
        d.fill.solid()
        d.fill.fore_color.rgb = color
        d.line.fill.background()
        return d

    def add_circle(self, slide, x, y, size, color=None):
        """Add a circle."""
        color = color or TELEKOM.MAGENTA
        c = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
        c.fill.solid()
        c.fill.fore_color.rgb = color
        c.line.fill.background()
        return c

    def add_line(self, slide, x, y, w, h, color=None):
        """Add a line (w≈0 for vertical, h≈0 for horizontal)."""
        color = color or TELEKOM.LIGHT_GRAY
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        line.fill.solid()
        line.fill.fore_color.rgb = color
        line.line.fill.background()
        return line

    def add_arrow(self, slide, x, y, w, h, color=None, direction='right'):
        """Add an arrow shape."""
        color = color or TELEKOM.MAGENTA
        shape_type = {
            'right': MSO_SHAPE.RIGHT_ARROW,
            'left': MSO_SHAPE.LEFT_ARROW,
            'up': MSO_SHAPE.UP_ARROW,
            'down': MSO_SHAPE.DOWN_ARROW,
        }.get(direction, MSO_SHAPE.RIGHT_ARROW)
        a = slide.shapes.add_shape(
            shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
        a.fill.solid()
        a.fill.fore_color.rgb = color
        a.line.fill.background()
        return a

    def add_accent_bar(self, slide, x, y, color=None,
                       direction='vertical', length=0.26, thickness=0.06):
        """Add a small accent bar (phase header marker, etc.)."""
        color = color or TELEKOM.MAGENTA
        if direction == 'vertical':
            return self.add_rect(slide, x, y, thickness, length, color)
        else:
            return self.add_rect(slide, x, y, length, thickness, color)

    # --- Deck-matched compound shapes ---

    def add_magenta_gradient_card(self, slide, x, y, w, h, text='',
                                  text_size=12, text_color=None, radius=16667):
        """Add a rounded rect with the signature DT Magenta gradient fill.
        Matches the gradient cards on slides 3, 6, 8 (Forschungsfragen, Blockers).
        """
        text_color = text_color or TELEKOM.WHITE
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.line.fill.background()
        _set_corner_radius(shape, radius)
        _apply_magenta_gradient(shape)
        if text:
            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(text_size)
            p.font.bold = True
            p.font.color.rgb = text_color
            p.alignment = PP_ALIGN.CENTER
        return shape

    def add_numbered_circle(self, slide, x, y, size, number,
                            gradient=True, color=None):
        """Add a numbered circle (like the flowchart connectors on slide 8).
        Uses flowChartConnector shape with Magenta gradient by default.
        """
        shape = slide.shapes.add_shape(
            MSO_SHAPE.FLOWCHART_CONNECTOR,
            Inches(x), Inches(y), Inches(size), Inches(size))
        if gradient:
            _apply_magenta_gradient(shape)
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = color or TELEKOM.MAGENTA
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = str(number)
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = TELEKOM.WHITE
        p.alignment = PP_ALIGN.CENTER
        return shape

    def add_info_card(self, slide, x, y, w, h, title='', body='',
                      bg_color=None):
        """Add an info card (like slides 10/11/13 — gray bg, title + body).
        Matches the EDEDED cards in the existing deck.
        """
        bg_color = bg_color or TELEKOM.CARD_BG
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.fill.background()
        _set_corner_radius(card, 16667)
        # Title label (inside card)
        if title:
            self.add_label(slide, x + 0.15, y + 0.1, w - 0.3, 0.3, title,
                           size=12, bold=True, color=TELEKOM.DARK)
        # Body text (inside card)
        if body:
            self.add_label(slide, x + 0.15, y + 0.4, w - 0.3, h - 0.5, body,
                           size=10, color=TELEKOM.DARK, wrap=True)
        return card

    def add_highlight_card(self, slide, x, y, w, h, text='',
                           bg_color=None, text_color=None):
        """Add a highlight card (like the Magenta E20074 cards on slides 11/13)."""
        bg_color = bg_color or TELEKOM.MAGENTA
        text_color = text_color or TELEKOM.WHITE
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.fill.background()
        _set_corner_radius(card, 16667)
        if text:
            tf = card.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.15)
            tf.margin_top = Inches(0.1)
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = text_color
        return card

    def add_process_flow(self, slide, items, x_start, y, item_w=2.3,
                         item_h=1.2, gap=0.3, highlight_indices=None):
        """Add a horizontal process flow (like slide 13 — boxes with arrows).
        items: list of (title, subtitle) tuples.
        highlight_indices: which items get Magenta fill (0-indexed).
        """
        highlight_indices = highlight_indices or []
        x = x_start
        for i, (title, subtitle) in enumerate(items):
            is_highlight = i in highlight_indices
            bg = TELEKOM.MAGENTA if is_highlight else TELEKOM.CARD_BG
            txt_color = TELEKOM.WHITE if is_highlight else TELEKOM.DARK
            self.add_info_card(slide, x, y, item_w, item_h, title, subtitle, bg)
            if is_highlight:
                # Override text colors for highlight cards
                pass  # info_card handles this internally with bg_color
            # Arrow between items
            if i < len(items) - 1:
                self.add_label(slide, x + item_w, y + item_h / 2 - 0.15,
                               gap, 0.3, '→', size=14, color=TELEKOM.GRAY,
                               align=PP_ALIGN.CENTER)
            x += item_w + gap

    # --- Glass morphism & Assembly-style primitives ---

    def add_glass_card(self, slide, x, y, w, h, text='', text_size=12,
                       text_color=None, alpha=20, fill_hex='FFFFFF',
                       radius=50000):
        """Add a glass morphism card — semi-transparent rounded rect.

        Extracted from Assembly deck: roundRect with #FFFFFF @ ~20% alpha.
        Great for overlaying on dark/image backgrounds.

        Args:
            alpha: transparency percentage (0=invisible, 100=opaque). Default 20%.
            fill_hex: fill color hex string without #. Default white.
            radius: corner radius (50000=full pill, 16667=standard). Default pill.
        """
        text_color = text_color or TELEKOM.WHITE
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.line.fill.background()
        _set_corner_radius(shape, radius)

        # Apply semi-transparent fill via XML (python-pptx has no alpha API)
        spPr = shape._element.find(qn('p:spPr'))
        if spPr is None:
            spPr = shape._element.find(qn('a:spPr'))
        # Remove existing fill
        for old_fill in spPr.findall(qn('a:solidFill')) + spPr.findall(qn('a:noFill')):
            spPr.remove(old_fill)
        solidFill = etree.SubElement(spPr, qn('a:solidFill'))
        srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'), attrib={'val': fill_hex})
        # Alpha: 0 = fully transparent, 100000 = fully opaque
        alpha_val = str(int(alpha * 1000))
        etree.SubElement(srgbClr, qn('a:alpha'), attrib={'val': alpha_val})

        if text:
            tf = shape.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.15)
            tf.margin_right = Inches(0.15)
            tf.margin_top = Inches(0.1)
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(text_size)
            p.font.color.rgb = text_color
            p.font.name = FONT_FALLBACK
        return shape

    def add_footer_bar(self, slide, number=None, section='', year='',
                       text_color=None, line_color=None, font_size=8):
        """Add a consistent footer bar at the bottom of a slide.

        Pattern from Assembly deck: number | section title | year
        with a thin horizontal line above.

        Args:
            number: slide number (int or str). None to omit.
            section: section/chapter title.
            year: year string (e.g. '2026').
        """
        text_color = text_color or TELEKOM.GRAY
        line_color = line_color or TELEKOM.LIGHT_GRAY
        sw = self.slide_width / 914400  # inches
        sh = self.slide_height / 914400

        # Horizontal line
        self.add_line(slide, 0.3, sh - 0.5, sw - 0.6, 0.01, line_color)

        # Number (left)
        if number is not None:
            self.add_label(slide, 0.3, sh - 0.42, 0.5, 0.3, str(number),
                           size=font_size, color=text_color, align=PP_ALIGN.LEFT)
        # Section (center)
        if section:
            self.add_label(slide, sw * 0.2, sh - 0.42, sw * 0.6, 0.3, section,
                           size=font_size, color=text_color, align=PP_ALIGN.CENTER)
        # Year (right)
        if year:
            self.add_label(slide, sw - 1.3, sh - 0.42, 1.0, 0.3, str(year),
                           size=font_size, color=text_color, align=PP_ALIGN.RIGHT)

    # --- Compound layout helpers ---

    def add_column_shading(self, slide, columns, x_start, y_start,
                           col_width, height, color=None):
        """Add alternating column background shading."""
        color = color or TELEKOM.SURFACE
        for i in range(columns):
            if i % 2 == 1:
                x = x_start + i * col_width
                self.add_rect(slide, x, y_start, col_width, height, color)

    def add_legend_item(self, slide, x, y, shape_type, color, label,
                        label_color=None):
        """Add a single legend entry (shape + label)."""
        label_color = label_color or TELEKOM.GRAY
        if shape_type == 'bar':
            self.add_rounded_bar(slide, x, y, 0.4, 0.14, color)
            self.add_label(slide, x + 0.5, y - 0.02, 0.8, 0.2,
                           label, size=8, color=label_color)
        elif shape_type == 'diamond':
            self.add_diamond(slide, x, y - 0.01, 0.15, color)
            self.add_label(slide, x + 0.22, y - 0.02, 0.8, 0.2,
                           label, size=8, color=label_color)
        elif shape_type == 'circle':
            self.add_circle(slide, x, y, 0.14, color)
            self.add_label(slide, x + 0.22, y - 0.02, 0.8, 0.2,
                           label, size=8, color=label_color)

    # --- Missing primitives: Tables, Images, Bullets, Notes ---

    def add_table(self, slide, x, y, w, h, data,
                  header=True, col_widths=None,
                  header_color=None, header_text_color=None,
                  stripe_color=None, font_size=10, header_font_size=11):
        """Add a styled table.

        Args:
            data: list of lists — first row is header if header=True.
            col_widths: list of widths in inches (auto-calculated if None).
            header_color: header row fill (default: MAGENTA).
            stripe_color: alternating row fill (default: SURFACE).
        """
        header_color = header_color or TELEKOM.MAGENTA
        header_text_color = header_text_color or TELEKOM.WHITE
        stripe_color = stripe_color or TELEKOM.SURFACE
        rows = len(data)
        cols = len(data[0]) if data else 0
        if rows == 0 or cols == 0:
            return None

        tbl_shape = slide.shapes.add_table(rows, cols,
                                           Inches(x), Inches(y),
                                           Inches(w), Inches(h))
        table = tbl_shape.table

        # Column widths
        if col_widths:
            for i, cw in enumerate(col_widths):
                table.columns[i].width = Inches(cw)
        else:
            col_w = w / cols
            for i in range(cols):
                table.columns[i].width = Inches(col_w)

        # Fill cells
        for row_idx, row_data in enumerate(data):
            for col_idx, cell_text in enumerate(row_data):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(cell_text)
                p = cell.text_frame.paragraphs[0]
                p.font.name = FONT_FALLBACK
                p.font.size = Pt(font_size)
                p.font.color.rgb = TELEKOM.DARK

                # Header styling
                if header and row_idx == 0:
                    p.font.bold = True
                    p.font.size = Pt(header_font_size)
                    p.font.color.rgb = header_text_color
                    self._set_cell_fill(cell, header_color)
                # Stripe
                elif row_idx % 2 == 0:
                    self._set_cell_fill(cell, stripe_color)

                # Vertical centering
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        return tbl_shape

    @staticmethod
    def _set_cell_fill(cell, color):
        """Set a table cell's fill color via XML."""
        tcPr = cell._tc.get_or_add_tcPr()
        # Remove existing fill
        for tag in ['a:solidFill', 'a:noFill']:
            el = tcPr.find(qn(tag))
            if el is not None:
                tcPr.remove(el)
        solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
        etree.SubElement(solidFill, qn('a:srgbClr'), attrib={
            'val': '%02X%02X%02X' % (color[0], color[1], color[2])
        })

    def add_image(self, slide, image_path, x, y, w=None, h=None):
        """Add an image to a slide. Specify w or h (or both) in inches.
        If only one dimension given, aspect ratio is preserved.
        """
        kwargs = {'left': Inches(x), 'top': Inches(y)}
        if w is not None:
            kwargs['width'] = Inches(w)
        if h is not None:
            kwargs['height'] = Inches(h)
        return slide.shapes.add_picture(str(image_path), **kwargs)

    def add_bullet_list(self, slide, x, y, w, h, items,
                        font_size=11, color=None, bullet_char='•',
                        line_spacing=1.2, indent_levels=None):
        """Add a bullet list. Each item is a string.

        Args:
            items: list of strings.
            indent_levels: optional list of ints (0=top, 1=sub, 2=sub-sub).
            bullet_char: bullet character (default '•').
        """
        color = color or TELEKOM.DARK
        indent_levels = indent_levels or [0] * len(items)
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True

        for i, text in enumerate(items):
            level = indent_levels[i] if i < len(indent_levels) else 0
            indent = '  ' * level
            prefix = bullet_char + ' ' if bullet_char else ''

            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.text = f"{indent}{prefix}{text}"
            p.font.size = Pt(font_size - level)  # Slightly smaller for sub-items
            p.font.color.rgb = color
            p.font.name = FONT_FALLBACK
            p.space_after = Pt(font_size * (line_spacing - 1))
            p.level = level

        return tb

    def add_rich_text(self, slide, x, y, w, h, paragraphs,
                      default_size=11, default_color=None, wrap=True):
        """Add multi-paragraph text with per-paragraph styling.

        Args:
            paragraphs: list of dicts with keys:
                text (str), size (int), bold (bool), italic (bool),
                color (RGBColor), align (PP_ALIGN).
                Only 'text' is required.
        """
        default_color = default_color or TELEKOM.DARK
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = wrap

        for i, para in enumerate(paragraphs):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.text = para.get('text', '')
            p.font.size = Pt(para.get('size', default_size))
            p.font.bold = para.get('bold', False)
            p.font.name = FONT_FALLBACK
            p.font.color.rgb = para.get('color', default_color)
            if 'align' in para:
                p.alignment = para['align']
            if para.get('italic'):
                p.font.italic = True

        return tb

    def set_speaker_notes(self, slide, text):
        """Set speaker notes for a slide."""
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = text

    def duplicate_slide(self, slide_index):
        """Duplicate a slide by index (0-based). Returns the new slide."""
        template = self.prs.slides[slide_index]
        # Copy slide layout
        layout = template.slide_layout
        new_slide = self.prs.slides.add_slide(layout)

        # Copy all shapes via XML
        for shape in template.shapes:
            el = deepcopy(shape._element)
            new_slide.shapes._spTree.append(el)

        # Remove default placeholders that came from layout
        for ph in list(new_slide.placeholders):
            # Keep only if it was in the original
            ph._element.getparent().remove(ph._element)

        return new_slide

    def find_and_replace(self, text_find, text_replace, slides=None):
        """Replace text across all slides (or a subset).

        Args:
            slides: optional list of slide indices (0-based). None = all.
        Returns: number of replacements made.
        """
        count = 0
        target_slides = (
            [self.prs.slides[i] for i in slides]
            if slides else self.prs.slides
        )
        for slide in target_slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if text_find in run.text:
                                run.text = run.text.replace(text_find, text_replace)
                                count += 1
                # Also check table cells
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            for para in cell.text_frame.paragraphs:
                                for run in para.runs:
                                    if text_find in run.text:
                                        run.text = run.text.replace(
                                            text_find, text_replace)
                                        count += 1
        return count


# ============================================================
# SLIDE RECIPES — High-level slide builders
# ============================================================

class SlideRecipes:
    """Pre-built slide templates for common presentation patterns.

    Usage:
        sb = SlideBuilder("deck.pptx")
        recipes = SlideRecipes(sb)
        recipes.title_slide("My Presentation", "Lars Boes — Feb 2026")
        recipes.section_divider("Chapter 1")
        recipes.kpi_dashboard([("Users", "1.2M", "+12%"), ...])
        sb.save()
    """

    # Slide dimensions (widescreen 13.33 x 7.5 inches)
    SLIDE_W = 13.333
    SLIDE_H = 7.5
    MARGIN = 0.6
    CONTENT_W = SLIDE_W - 2 * MARGIN

    def __init__(self, sb: SlideBuilder):
        self.sb = sb

    def title_slide(self, title, subtitle='', notes=''):
        """Full-width title slide with centered text."""
        slide = self.sb.add_blank_slide()
        # Magenta accent bar at top
        self.sb.add_rect(slide, 0, 0, self.SLIDE_W, 0.08, TELEKOM.MAGENTA)
        # Title — centered
        self.sb.add_label(slide, self.MARGIN, 2.2, self.CONTENT_W, 1.5,
                          title, size=40, bold=True, color=TELEKOM.DARK,
                          align=PP_ALIGN.CENTER)
        # Subtitle
        if subtitle:
            self.sb.add_label(slide, self.MARGIN, 3.8, self.CONTENT_W, 0.6,
                              subtitle, size=16, color=TELEKOM.GRAY,
                              align=PP_ALIGN.CENTER)
        if notes:
            self.sb.set_speaker_notes(slide, notes)
        return slide

    def section_divider(self, title, subtitle='', notes=''):
        """Section divider / chapter separator — Magenta left bar + large title."""
        slide = self.sb.add_blank_slide()
        # Full-height magenta accent bar (left)
        self.sb.add_rect(slide, 0, 0, 0.15, self.SLIDE_H, TELEKOM.MAGENTA)
        # Title
        self.sb.add_label(slide, 1.0, 2.5, 10, 1.2, title,
                          size=36, bold=True, color=TELEKOM.DARK)
        if subtitle:
            self.sb.add_label(slide, 1.0, 3.8, 10, 0.6, subtitle,
                              size=14, color=TELEKOM.GRAY)
        if notes:
            self.sb.set_speaker_notes(slide, notes)
        return slide

    def content_slide(self, title, bullets, subtitle='', notes=''):
        """Standard content slide — title + bullet list."""
        slide = self.sb.add_blank_slide()
        # Magenta top accent
        self.sb.add_rect(slide, 0, 0, self.SLIDE_W, 0.04, TELEKOM.MAGENTA)
        # Title
        self.sb.add_label(slide, self.MARGIN, 0.4, self.CONTENT_W, 0.6,
                          title, size=28, bold=True, color=TELEKOM.DARK)
        if subtitle:
            self.sb.add_label(slide, self.MARGIN, 1.0, self.CONTENT_W, 0.4,
                              subtitle, size=12, color=TELEKOM.GRAY)
        # Bullet list
        top = 1.5 if subtitle else 1.2
        self.sb.add_bullet_list(slide, self.MARGIN, top,
                                self.CONTENT_W, self.SLIDE_H - top - 0.8,
                                bullets, font_size=14)
        if notes:
            self.sb.set_speaker_notes(slide, notes)
        return slide

    def kpi_dashboard(self, metrics, title='Key Metrics', notes=''):
        """KPI dashboard — up to 4 metric cards in a row.

        Args:
            metrics: list of tuples (label, value, delta).
                     delta is optional (e.g., "+12%", "−3%").
        """
        slide = self.sb.add_blank_slide()
        self.sb.add_rect(slide, 0, 0, self.SLIDE_W, 0.04, TELEKOM.MAGENTA)
        self.sb.add_label(slide, self.MARGIN, 0.4, self.CONTENT_W, 0.6,
                          title, size=28, bold=True, color=TELEKOM.DARK)

        n = min(len(metrics), 4)
        card_gap = 0.3
        total_gap = card_gap * (n - 1)
        card_w = (self.CONTENT_W - total_gap) / n
        card_h = 2.5
        card_y = 2.5

        for i, metric in enumerate(metrics):
            label = metric[0]
            value = metric[1]
            delta = metric[2] if len(metric) > 2 else None

            card_x = self.MARGIN + i * (card_w + card_gap)

            # Card background
            self.sb.add_rounded_rect(slide, card_x, card_y, card_w, card_h,
                                     TELEKOM.CARD_BG)

            # Value — large centered
            self.sb.add_label(slide, card_x, card_y + 0.4, card_w, 1.0,
                              str(value), size=36, bold=True,
                              color=TELEKOM.MAGENTA, align=PP_ALIGN.CENTER)
            # Label
            self.sb.add_label(slide, card_x, card_y + 1.5, card_w, 0.4,
                              label, size=12, color=TELEKOM.DARK,
                              align=PP_ALIGN.CENTER)
            # Delta
            if delta:
                delta_str = str(delta)
                delta_color = (TELEKOM.GREEN if delta_str.startswith('+')
                               else TELEKOM.RED if delta_str.startswith('-') or delta_str.startswith('−')
                               else TELEKOM.GRAY)
                self.sb.add_label(slide, card_x, card_y + 1.9, card_w, 0.3,
                                  delta_str, size=11, bold=True,
                                  color=delta_color, align=PP_ALIGN.CENTER)

        if notes:
            self.sb.set_speaker_notes(slide, notes)
        return slide

    def comparison_slide(self, title, left_title, left_items,
                         right_title, right_items,
                         left_color=None, right_color=None, notes=''):
        """Two-column comparison (e.g., Pros/Cons, Before/After, Option A/B)."""
        left_color = left_color or TELEKOM.TEAL
        right_color = right_color or TELEKOM.MAGENTA
        slide = self.sb.add_blank_slide()
        self.sb.add_rect(slide, 0, 0, self.SLIDE_W, 0.04, TELEKOM.MAGENTA)
        self.sb.add_label(slide, self.MARGIN, 0.4, self.CONTENT_W, 0.6,
                          title, size=28, bold=True, color=TELEKOM.DARK)

        col_w = (self.CONTENT_W - 0.4) / 2
        col_h = 4.5
        top_y = 1.5

        for i, (col_title, items, color) in enumerate([
            (left_title, left_items, left_color),
            (right_title, right_items, right_color),
        ]):
            col_x = self.MARGIN + i * (col_w + 0.4)
            # Column header bar
            self.sb.add_rounded_rect(slide, col_x, top_y, col_w, 0.5, color)
            self.sb.add_label(slide, col_x + 0.2, top_y + 0.08, col_w - 0.4, 0.4,
                              col_title, size=14, bold=True, color=TELEKOM.WHITE)
            # Items
            self.sb.add_bullet_list(slide, col_x + 0.2, top_y + 0.7,
                                    col_w - 0.4, col_h - 0.7,
                                    items, font_size=12)

        if notes:
            self.sb.set_speaker_notes(slide, notes)
        return slide

    def agenda_slide(self, items, current_index=None, title='Agenda', notes=''):
        """Agenda / table of contents with optional highlight of current section.

        Args:
            items: list of section titles.
            current_index: 0-based index to highlight (None = no highlight).
        """
        slide = self.sb.add_blank_slide()
        self.sb.add_rect(slide, 0, 0, self.SLIDE_W, 0.04, TELEKOM.MAGENTA)
        self.sb.add_label(slide, self.MARGIN, 0.4, self.CONTENT_W, 0.6,
                          title, size=28, bold=True, color=TELEKOM.DARK)

        item_h = 0.6
        item_gap = 0.15
        start_y = 1.8
        item_w = 8.0
        item_x = (self.SLIDE_W - item_w) / 2

        for i, item in enumerate(items):
            y = start_y + i * (item_h + item_gap)
            is_current = (i == current_index)

            bg = TELEKOM.MAGENTA if is_current else TELEKOM.CARD_BG
            text_color = TELEKOM.WHITE if is_current else TELEKOM.DARK

            self.sb.add_rounded_rect(slide, item_x, y, item_w, item_h, bg)
            # Number circle
            self.sb.add_numbered_circle(
                slide, item_x + 0.15, y + 0.08, 0.42, str(i + 1),
                gradient=is_current, color=TELEKOM.GRAY if not is_current else None)
            # Text
            self.sb.add_label(slide, item_x + 0.75, y + 0.12, item_w - 1.0, 0.35,
                              item, size=14, bold=is_current, color=text_color)

        if notes:
            self.sb.set_speaker_notes(slide, notes)
        return slide

    def image_slide(self, title, image_path, caption='',
                    image_side='right', notes=''):
        """Content slide with image — text left/right, image on the other side.

        Args:
            image_side: 'right' (text left, image right) or 'left'.
        """
        slide = self.sb.add_blank_slide()
        self.sb.add_rect(slide, 0, 0, self.SLIDE_W, 0.04, TELEKOM.MAGENTA)

        if image_side == 'right':
            text_x, img_x = self.MARGIN, 7.0
        else:
            text_x, img_x = 6.5, self.MARGIN

        # Title
        self.sb.add_label(slide, text_x, 0.4, 6.0, 0.6,
                          title, size=28, bold=True, color=TELEKOM.DARK)
        # Caption
        if caption:
            self.sb.add_label(slide, text_x, 1.2, 6.0, 4.5,
                              caption, size=13, color=TELEKOM.DARK, wrap=True)
        # Image
        self.sb.add_image(slide, image_path, img_x, 1.0, w=5.5)

        if notes:
            self.sb.set_speaker_notes(slide, notes)
        return slide

    def table_slide(self, title, data, col_widths=None, notes=''):
        """Slide with a styled data table.

        Args:
            data: list of lists — first row is header.
            col_widths: optional list of column widths in inches.
        """
        slide = self.sb.add_blank_slide()
        self.sb.add_rect(slide, 0, 0, self.SLIDE_W, 0.04, TELEKOM.MAGENTA)
        self.sb.add_label(slide, self.MARGIN, 0.4, self.CONTENT_W, 0.6,
                          title, size=28, bold=True, color=TELEKOM.DARK)

        rows = len(data)
        row_h = min(0.5, 5.0 / max(rows, 1))
        table_h = rows * row_h
        table_w = self.CONTENT_W

        self.sb.add_table(slide, self.MARGIN, 1.5, table_w, table_h,
                          data, col_widths=col_widths)

        if notes:
            self.sb.set_speaker_notes(slide, notes)
        return slide

    def chart_slide(self, title, chart_type, categories, series_data,
                    notes='', **chart_kwargs):
        """Slide with a titled chart. Wraps ChartBuilder methods.

        Args:
            chart_type: 'column', 'bar', 'line', 'pie', 'doughnut', 'scatter'.
            categories: list of category labels.
            series_data: for pie/doughnut: list of values.
                         for scatter: list of (name, [(x,y),...]) tuples.
                         for others: list of (series_name, values) tuples.
            chart_kwargs: extra args passed to the ChartBuilder method.
        """
        slide = self.sb.add_blank_slide()
        self.sb.add_rect(slide, 0, 0, self.SLIDE_W, 0.04, TELEKOM.MAGENTA)

        charts = ChartBuilder(self.sb)
        method = getattr(charts, f'{chart_type}_chart', None)
        if method is None:
            raise ValueError(f"Unknown chart type '{chart_type}'. "
                             f"Use: column, bar, line, pie, doughnut, scatter")

        if chart_type in ('pie', 'doughnut'):
            method(slide, title, categories, series_data,
                   x=self.MARGIN, y=0.8, w=self.CONTENT_W, h=6.2,
                   **chart_kwargs)
        elif chart_type == 'scatter':
            method(slide, title, series_data,
                   x=self.MARGIN, y=0.8, w=self.CONTENT_W, h=6.2,
                   **chart_kwargs)
        else:
            method(slide, title, categories, series_data,
                   x=self.MARGIN, y=0.8, w=self.CONTENT_W, h=6.2,
                   **chart_kwargs)

        if notes:
            self.sb.set_speaker_notes(slide, notes)
        return slide

    # --- Assembly-inspired recipes ---

    def numbered_section_divider(self, title, number, subtitle='',
                                 circle_color=None, notes=''):
        """Section divider with large numbered circle (Assembly deck pattern).

        Creates a chapter break slide with a magenta numbered circle,
        large title, and optional subtitle in a glass pill.

        Args:
            number: chapter/section number (displayed in circle).
            subtitle: optional tagline shown in glass pill below title.
        """
        circle_color = circle_color or TELEKOM.MAGENTA
        slide = self.sb.add_blank_slide()

        # Magenta accent bar at top
        self.sb.add_rect(slide, 0, 0, self.SLIDE_W, 0.04, TELEKOM.MAGENTA)

        # Numbered circle (upper-left area)
        circle_size = 1.3
        c = self.sb.add_circle(slide, self.MARGIN + 0.2, 1.0, circle_size, circle_color)
        tf = c.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = str(number).zfill(2)
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = TELEKOM.WHITE
        p.font.name = FONT_FALLBACK
        p.alignment = PP_ALIGN.CENTER

        # Large title
        self.sb.add_label(slide, self.MARGIN, 2.8, self.CONTENT_W, 1.5,
                          title, size=44, bold=True, color=TELEKOM.DARK)

        # Subtitle glass pill
        if subtitle:
            self.sb.add_glass_card(slide, self.MARGIN, 4.8,
                                   self.CONTENT_W * 0.6, 0.6,
                                   subtitle, text_size=12,
                                   text_color=TELEKOM.GRAY,
                                   alpha=40, fill_hex='EDEDED', radius=50000)

        if notes:
            self.sb.set_speaker_notes(slide, notes)
        return slide

    def icon_grid_slide(self, title, items, cols=3, icon_color=None, notes=''):
        """Grid of icon circles with titles and descriptions.

        Pattern from Assembly slide 6: icon circle + bold title + description.

        Args:
            items: list of (icon_label, title, description) tuples.
                   icon_label is a short string (emoji, letter, or 1-2 chars)
                   displayed inside the circle.
            cols: number of columns (2, 3, or 4).
        """
        icon_color = icon_color or TELEKOM.MAGENTA
        slide = self.sb.add_blank_slide()
        self.sb.add_rect(slide, 0, 0, self.SLIDE_W, 0.04, TELEKOM.MAGENTA)
        self.sb.add_label(slide, self.MARGIN, 0.4, self.CONTENT_W, 0.6,
                          title, size=28, bold=True, color=TELEKOM.DARK)

        col_w = self.CONTENT_W / cols
        circle_size = 0.7
        y_start = 1.5
        row_h = 2.8

        for i, item in enumerate(items):
            icon_label, item_title, description = item
            col = i % cols
            row = i // cols
            x = self.MARGIN + col * col_w + (col_w - circle_size) / 2
            y = y_start + row * row_h

            # Alternating circle colors
            is_alt = i % 2 == 1
            c_color = TELEKOM.SURFACE if is_alt else icon_color
            t_color = TELEKOM.DARK if is_alt else TELEKOM.WHITE

            # Circle with icon
            c = self.sb.add_circle(slide, x, y, circle_size, c_color)
            tf = c.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.text = str(icon_label)
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = t_color
            p.font.name = FONT_FALLBACK
            p.alignment = PP_ALIGN.CENTER

            # Title below circle
            text_x = self.MARGIN + col * col_w
            self.sb.add_label(slide, text_x, y + circle_size + 0.15,
                              col_w, 0.35, item_title,
                              size=14, bold=True, color=TELEKOM.DARK,
                              align=PP_ALIGN.CENTER)

            # Description
            self.sb.add_label(slide, text_x, y + circle_size + 0.5,
                              col_w, 0.8, description,
                              size=10, color=TELEKOM.GRAY,
                              align=PP_ALIGN.CENTER, wrap=True)

        if notes:
            self.sb.set_speaker_notes(slide, notes)
        return slide

    def milestone_timeline_slide(self, title, milestones, notes=''):
        """Vertical cascading milestone timeline (Assembly slide 53 pattern).

        Creates a staircase/cascade of numbered milestones with description
        cards, connected by a vertical line element.

        Args:
            milestones: list of (number_or_label, bullets: list[str]) tuples.
                        Each becomes a numbered circle + glass card.
        """
        slide = self.sb.add_blank_slide()
        self.sb.add_rect(slide, 0, 0, self.SLIDE_W, 0.04, TELEKOM.MAGENTA)
        self.sb.add_label(slide, self.MARGIN, 0.4, self.CONTENT_W, 0.6,
                          title, size=28, bold=True, color=TELEKOM.DARK)

        n = len(milestones)
        available_h = 5.5  # Space for milestones
        step_h = min(available_h / max(n, 1), 1.8)
        circle_size = 0.55
        card_w = 8.0
        x_indent_step = 0.6  # Each milestone indents further right

        y = 1.5
        for i, (label, bullets) in enumerate(milestones):
            x_base = self.MARGIN + i * x_indent_step

            # Numbered circle
            c = self.sb.add_circle(slide, x_base, y, circle_size, TELEKOM.MAGENTA)
            tf = c.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.text = str(label)
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = TELEKOM.WHITE
            p.font.name = FONT_FALLBACK
            p.alignment = PP_ALIGN.CENTER

            # Connector line to next milestone
            if i < n - 1:
                line_x = x_base + circle_size / 2 - 0.01
                self.sb.add_line(slide, line_x, y + circle_size,
                                 0.02, step_h - circle_size, TELEKOM.MAGENTA)

            # Glass card with bullets
            card_x = x_base + circle_size + 0.2
            remaining_w = self.SLIDE_W - card_x - self.MARGIN
            actual_card_w = min(card_w, remaining_w)
            card_h = max(0.6, len(bullets) * 0.28 + 0.2)

            card = self.sb.add_glass_card(slide, card_x, y,
                                          actual_card_w, card_h,
                                          alpha=30, fill_hex='EDEDED')
            # Add bullet text
            self.sb.add_bullet_list(slide, card_x + 0.15, y + 0.08,
                                    actual_card_w - 0.3, card_h - 0.15,
                                    bullets, font_size=10, bullet_char='•')

            y += step_h

        if notes:
            self.sb.set_speaker_notes(slide, notes)
        return slide

    def risk_grid_slide(self, title, risks, cols=2, notes=''):
        """Grid of risk/info cards with numbered circles (Assembly slide 50).

        Args:
            risks: list of (id_label, bullets: list[str]) tuples.
                   id_label: short identifier (e.g. "R1", "R2").
            cols: 2 or 3 columns.
        """
        slide = self.sb.add_blank_slide()
        self.sb.add_rect(slide, 0, 0, self.SLIDE_W, 0.04, TELEKOM.MAGENTA)
        self.sb.add_label(slide, self.MARGIN, 0.4, self.CONTENT_W, 0.6,
                          title, size=28, bold=True, color=TELEKOM.DARK)

        col_w = (self.CONTENT_W - 0.4) / cols
        card_h = 1.6
        row_gap = 0.3
        col_gap = 0.4
        circle_size = 0.45
        y_start = 1.5

        for i, (label, bullets) in enumerate(risks):
            col = i % cols
            row = i // cols
            x = self.MARGIN + col * (col_w + col_gap)
            y = y_start + row * (card_h + row_gap)

            # Numbered circle
            c = self.sb.add_circle(slide, x, y, circle_size, TELEKOM.MAGENTA)
            tf = c.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.text = str(label)
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = TELEKOM.WHITE
            p.font.name = FONT_FALLBACK
            p.alignment = PP_ALIGN.CENTER

            # Info card next to circle
            card_x = x + circle_size + 0.15
            actual_card_w = col_w - circle_size - 0.15
            self.sb.add_glass_card(slide, card_x, y,
                                   actual_card_w, card_h,
                                   alpha=30, fill_hex='EDEDED')
            self.sb.add_bullet_list(slide, card_x + 0.1, y + 0.08,
                                    actual_card_w - 0.2, card_h - 0.15,
                                    bullets, font_size=10, bullet_char='•')

        if notes:
            self.sb.set_speaker_notes(slide, notes)
        return slide

    def phase_tree_slide(self, title, phase_number, phase_name, items,
                         notes=''):
        """Phase breakdown with tree connector (Assembly slide 42-45 pattern).

        Shows a numbered phase circle with vertical/horizontal branch lines
        leading to item cards, each with an arrow pointing to description bullets.

        Args:
            phase_number: number displayed in circle.
            phase_name: phase title next to circle.
            items: list of (item_name, bullets: list[str]) tuples.
        """
        slide = self.sb.add_blank_slide()
        self.sb.add_rect(slide, 0, 0, self.SLIDE_W, 0.04, TELEKOM.MAGENTA)
        self.sb.add_label(slide, self.MARGIN, 0.4, self.CONTENT_W, 0.6,
                          title, size=28, bold=True, color=TELEKOM.DARK)

        # Phase circle + title
        circle_size = 0.7
        circle_x = self.MARGIN + 0.3
        circle_y = 1.5

        c = self.sb.add_circle(slide, circle_x, circle_y, circle_size, TELEKOM.MAGENTA)
        tf = c.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = str(phase_number)
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = TELEKOM.WHITE
        p.font.name = FONT_FALLBACK
        p.alignment = PP_ALIGN.CENTER

        self.sb.add_label(slide, circle_x + circle_size + 0.2, circle_y + 0.1,
                          4, 0.5, phase_name,
                          size=22, bold=True, color=TELEKOM.DARK)

        # Tree items
        n = len(items)
        item_h = 0.8
        item_gap = 0.25
        total_h = n * item_h + (n - 1) * item_gap
        tree_x = circle_x + circle_size / 2  # Vertical line x
        y_start = circle_y + circle_size + 0.3

        # Vertical trunk line
        self.sb.add_line(slide, tree_x - 0.01, circle_y + circle_size,
                         0.02, total_h + 0.3, TELEKOM.MAGENTA)

        card_w = 3.0
        card_x = tree_x + 0.5
        arrow_x = card_x + card_w + 0.15
        bullet_x = arrow_x + 0.6
        bullet_w = self.SLIDE_W - bullet_x - self.MARGIN

        for i, (item_name, bullets) in enumerate(items):
            y = y_start + i * (item_h + item_gap)

            # Horizontal branch
            self.sb.add_line(slide, tree_x, y + item_h / 2 - 0.01,
                             0.5, 0.02, TELEKOM.MAGENTA)

            # Item card (gray)
            self.sb.add_glass_card(slide, card_x, y, card_w, item_h,
                                   item_name, text_size=11,
                                   text_color=TELEKOM.DARK,
                                   alpha=40, fill_hex='EDEDED')

            # Arrow
            self.sb.add_arrow(slide, arrow_x, y + item_h / 2 - 0.15,
                              0.45, 0.3, TELEKOM.MAGENTA)

            # Bullets
            self.sb.add_bullet_list(slide, bullet_x, y + 0.05,
                                    bullet_w, item_h - 0.1,
                                    bullets, font_size=10, bullet_char='•')

        if notes:
            self.sb.set_speaker_notes(slide, notes)
        return slide

    def team_slide(self, title, members, notes=''):
        """Team portrait slide (Assembly slide 67 pattern).

        Circular photos with name and optional quote underneath.

        Args:
            members: list of (name, image_path, quote) tuples.
                     image_path can be None (shows colored circle placeholder).
                     quote can be empty string.
        """
        slide = self.sb.add_blank_slide()
        self.sb.add_rect(slide, 0, 0, self.SLIDE_W, 0.04, TELEKOM.MAGENTA)
        self.sb.add_label(slide, self.MARGIN, 0.4, self.CONTENT_W, 0.6,
                          title, size=28, bold=True, color=TELEKOM.DARK)

        n = len(members)
        col_w = self.CONTENT_W / max(n, 1)
        photo_size = min(2.5, col_w - 0.4)
        y_photo = 1.8

        for i, member in enumerate(members):
            name = member[0]
            image_path = member[1] if len(member) > 1 else None
            quote = member[2] if len(member) > 2 else ''

            cx = self.MARGIN + i * col_w + (col_w - photo_size) / 2

            if image_path and Path(image_path).exists():
                self.sb.add_image(slide, image_path, cx, y_photo,
                                  w=photo_size, h=photo_size)
            else:
                # Placeholder circle
                colors = [TELEKOM.MAGENTA, TELEKOM.TEAL, TELEKOM.PURPLE,
                          TELEKOM.SKY_BLUE]
                c = self.sb.add_circle(slide, cx, y_photo, photo_size,
                                       colors[i % len(colors)])
                tf = c.text_frame
                p = tf.paragraphs[0]
                p.text = name[0].upper()
                p.font.size = Pt(36)
                p.font.bold = True
                p.font.color.rgb = TELEKOM.WHITE
                p.alignment = PP_ALIGN.CENTER

            # Name
            text_x = self.MARGIN + i * col_w
            name_y = y_photo + photo_size + 0.3
            self.sb.add_label(slide, text_x, name_y, col_w, 0.4,
                              name, size=16, bold=True, color=TELEKOM.DARK,
                              align=PP_ALIGN.CENTER)

            # Quote
            if quote:
                self.sb.add_label(slide, text_x + 0.1, name_y + 0.5,
                                  col_w - 0.2, 1.0,
                                  f'"{quote}"',
                                  size=10, color=TELEKOM.GRAY,
                                  align=PP_ALIGN.CENTER, wrap=True)

        if notes:
            self.sb.set_speaker_notes(slide, notes)
        return slide


class GanttBuilder:
    """Build Gantt-style project overview slides.

    Usage:
        sb = SlideBuilder("deck.pptx")
        slide = sb.add_blank_slide()
        gantt = GanttBuilder(sb, slide,
            title="Projektübersicht",
            subtitle="My Project — 2026",
            months=["Mär", "Apr", "Mai", "Jun", "Jul"])
        gantt.add_phase("Phase 0: Setup")
        gantt.add_task("Task A", 0, 0.0, 0, 0.5, TELEKOM.MAGENTA)
        gantt.add_milestone("Key Date", 0, 0.5)
        gantt.add_today_marker(-0.15)
        gantt.add_legend([("Erledigt", "bar", TELEKOM.GREEN), ...])
    """

    def __init__(self, sb: SlideBuilder, slide,
                 title: str, subtitle: str,
                 months: list[str],
                 timeline_left: float = 3.3,
                 timeline_right: float = 12.8,
                 top_header: float = 1.15,
                 top_rows: float = 1.45,
                 row_spacing: float = 0.30,
                 bar_height: float = 0.18,
                 left_margin: float = 0.4,
                 label_width: float = 2.8):
        self.sb = sb
        self.slide = slide
        self.months = months
        self.tl_left = timeline_left
        self.tl_right = timeline_right
        self.tl_width = timeline_right - timeline_left
        self.top_header = top_header
        self.top_rows = top_rows
        self.row_spacing = row_spacing
        self.bar_h = bar_height
        self.left_margin = left_margin
        self.label_w = label_width
        self._row = 0

        sb.add_label(slide, 0.5, 0.25, 6, 0.6, title,
                     size=32, bold=True, color=TELEKOM.DARK)
        sb.add_label(slide, 0.5, 0.8, 6.5, 0.3, subtitle,
                     size=11, color=TELEKOM.GRAY)

        col_w = self.tl_width / len(months)
        for i, m in enumerate(months):
            x = self._month_x(i)
            if i % 2 == 1:
                sb.add_rect(slide, x, top_rows, col_w,
                            row_spacing * 25, TELEKOM.SURFACE)
            sb.add_label(slide, x, top_header, col_w, 0.3, m,
                         size=11, bold=True, color=TELEKOM.DARK,
                         align=PP_ALIGN.CENTER)

    def _month_x(self, month_idx, fraction=0.0):
        return self.tl_left + (month_idx + fraction) / len(self.months) * self.tl_width

    def _current_y(self):
        return self.top_rows + self._row * self.row_spacing

    def add_phase(self, label: str):
        y = self._current_y()
        self.sb.add_accent_bar(self.slide, self.left_margin, y + 0.02,
                               length=0.26, thickness=0.06)
        self.sb.add_label(self.slide, self.left_margin + 0.15, y,
                          self.label_w, 0.3, label,
                          size=11, bold=True, color=TELEKOM.DARK)
        self._row += 1

    def add_task(self, label, start_month, start_frac, end_month, end_frac,
                 color=None):
        color = color or TELEKOM.MAGENTA
        y = self._current_y()
        self.sb.add_label(self.slide, self.left_margin + 0.15, y,
                          self.label_w, 0.22, label,
                          size=9, color=TELEKOM.DARK)
        x1 = self._month_x(start_month, start_frac)
        x2 = self._month_x(end_month, end_frac)
        bar_w = max(x2 - x1, 0.25)
        self.sb.add_rounded_bar(self.slide, x1, y + 0.03, bar_w, self.bar_h, color)
        self._row += 1

    def add_milestone(self, label, month, fraction, color=None):
        color = color or TELEKOM.MAGENTA
        y = self._current_y()
        self.sb.add_label(self.slide, self.left_margin + 0.3, y,
                          self.label_w, 0.22, "🎯 " + label,
                          size=8.5, color=TELEKOM.MAGENTA)
        mx = self._month_x(month, fraction)
        self.sb.add_diamond(self.slide, mx - 0.07, y + 0.02, 0.15, color)
        self._row += 1

    def add_today_marker(self, month_fraction=0.0, month_index=0,
                         label="📍 Heute"):
        x = self._month_x(month_index, month_fraction)
        top = self.top_rows - 0.05
        height = self._row * self.row_spacing + 0.1
        self.sb.add_rect(self.slide, x, top, 0.02, height, TELEKOM.RED)
        self.sb.add_label(self.slide, x - 0.3, top - 0.25, 0.7, 0.2,
                          label, size=9, bold=True, color=TELEKOM.RED,
                          align=PP_ALIGN.CENTER)

    def add_legend(self, items, y=None, x_start=7.0):
        y = y or 0.85
        x = x_start
        for label, shape_type, color in items:
            self.sb.add_legend_item(self.slide, x, y, shape_type, color, label)
            x += 1.3 + len(label) * 0.05

    def add_slide_number(self, slide_num=None):
        num = slide_num or len(self.sb.prs.slides)
        self.sb.add_label(self.slide, 12.0, 7.1, 1, 0.3, str(num),
                          size=10, color=TELEKOM.GRAY, align=PP_ALIGN.RIGHT)


# ============================================================
# CHART BUILDER
# ============================================================

class ChartBuilder:
    """Native PowerPoint chart generation via python-pptx.

    Builds editable charts (not images) — users can modify data in PowerPoint.
    Theme colors from the deck are applied automatically.

    Usage:
        sb = SlideBuilder("deck.pptx")
        charts = ChartBuilder(sb)

        charts.column_chart(slide, "Q1 Revenue", ["East", "West", "Mid"],
                            [("2025", [19.2, 21.4, 16.7]),
                             ("2026", [22.3, 28.6, 15.2])])

        charts.pie_chart(slide, "Market Share",
                         ["Mobile", "Fixed", "TV", "Other"],
                         [0.45, 0.30, 0.15, 0.10])

        charts.line_chart(slide, "Trend", ["Jan", "Feb", "Mar"],
                          [("Users", [100, 150, 200])])
    """

    def __init__(self, sb: SlideBuilder):
        self.sb = sb

    @staticmethod
    def _apply_chart_style(chart, has_legend=True, font_size=10):
        """Apply consistent styling to any chart."""
        chart.font.size = Pt(font_size)
        if has_legend:
            chart.has_legend = True
            chart.legend.include_in_layout = False
            chart.legend.font.size = Pt(font_size)

    def column_chart(self, slide, title, categories, series_data,
                     x=0.6, y=1.5, w=12.0, h=5.5,
                     stacked=False, has_legend=True, data_labels=False):
        """Add a column (vertical bar) chart.

        Args:
            categories: list of category names.
            series_data: list of (series_name, values) tuples.
            stacked: if True, use stacked columns.
        """
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE

        chart_data = CategoryChartData()
        chart_data.categories = categories
        for name, values in series_data:
            chart_data.add_series(name, values)

        chart_type = (XL_CHART_TYPE.COLUMN_STACKED if stacked
                      else XL_CHART_TYPE.COLUMN_CLUSTERED)

        graphic_frame = slide.shapes.add_chart(
            chart_type, Inches(x), Inches(y), Inches(w), Inches(h), chart_data)
        chart = graphic_frame.chart
        self._apply_chart_style(chart, has_legend)

        if data_labels:
            chart.plots[0].has_data_labels = True
            chart.plots[0].data_labels.font.size = Pt(9)

        if title:
            chart.has_title = True
            chart.chart_title.text_frame.paragraphs[0].text = title
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)

        return chart

    def bar_chart(self, slide, title, categories, series_data,
                  x=0.6, y=1.5, w=12.0, h=5.5,
                  stacked=False, has_legend=True, data_labels=False):
        """Add a horizontal bar chart.

        Args:
            categories: list of category names.
            series_data: list of (series_name, values) tuples.
        """
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE

        chart_data = CategoryChartData()
        chart_data.categories = categories
        for name, values in series_data:
            chart_data.add_series(name, values)

        chart_type = (XL_CHART_TYPE.BAR_STACKED if stacked
                      else XL_CHART_TYPE.BAR_CLUSTERED)

        graphic_frame = slide.shapes.add_chart(
            chart_type, Inches(x), Inches(y), Inches(w), Inches(h), chart_data)
        chart = graphic_frame.chart
        self._apply_chart_style(chart, has_legend)

        if data_labels:
            chart.plots[0].has_data_labels = True

        if title:
            chart.has_title = True
            chart.chart_title.text_frame.paragraphs[0].text = title
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)

        return chart

    def line_chart(self, slide, title, categories, series_data,
                   x=0.6, y=1.5, w=12.0, h=5.5,
                   smooth=False, has_legend=True, data_labels=False):
        """Add a line chart.

        Args:
            categories: list of category names (x-axis labels).
            series_data: list of (series_name, values) tuples.
            smooth: if True, smooth all series lines into curves.
        """
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE

        chart_data = CategoryChartData()
        chart_data.categories = categories
        for name, values in series_data:
            chart_data.add_series(name, values)

        graphic_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE, Inches(x), Inches(y),
            Inches(w), Inches(h), chart_data)
        chart = graphic_frame.chart
        self._apply_chart_style(chart, has_legend)

        if smooth:
            for series in chart.series:
                series.smooth = True

        if data_labels:
            chart.plots[0].has_data_labels = True

        if title:
            chart.has_title = True
            chart.chart_title.text_frame.paragraphs[0].text = title
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)

        return chart

    def pie_chart(self, slide, title, categories, values,
                  x=0.6, y=1.5, w=12.0, h=5.5,
                  data_labels=True, number_format='0%'):
        """Add a pie chart.

        Args:
            categories: list of slice labels.
            values: list of numeric values (will be shown as percentages).
        """
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION

        chart_data = CategoryChartData()
        chart_data.categories = categories
        chart_data.add_series('Series 1', values)

        graphic_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.PIE, Inches(x), Inches(y),
            Inches(w), Inches(h), chart_data)
        chart = graphic_frame.chart
        self._apply_chart_style(chart, has_legend=True)

        from pptx.enum.chart import XL_LEGEND_POSITION
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM

        if data_labels:
            chart.plots[0].has_data_labels = True
            labels = chart.plots[0].data_labels
            labels.number_format = number_format
            labels.font.size = Pt(10)
            labels.position = XL_LABEL_POSITION.OUTSIDE_END

        if title:
            chart.has_title = True
            chart.chart_title.text_frame.paragraphs[0].text = title
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)

        return chart

    def scatter_chart(self, slide, title, series_data,
                      x=0.6, y=1.5, w=12.0, h=5.5,
                      has_legend=True):
        """Add an XY scatter chart.

        Args:
            series_data: list of (series_name, [(x, y), ...]) tuples.
        """
        from pptx.chart.data import XyChartData
        from pptx.enum.chart import XL_CHART_TYPE

        chart_data = XyChartData()
        for name, points in series_data:
            series = chart_data.add_series(name)
            for px, py in points:
                series.add_data_point(px, py)

        graphic_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.XY_SCATTER, Inches(x), Inches(y),
            Inches(w), Inches(h), chart_data)
        chart = graphic_frame.chart
        self._apply_chart_style(chart, has_legend)

        if title:
            chart.has_title = True
            chart.chart_title.text_frame.paragraphs[0].text = title
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)

        return chart

    def doughnut_chart(self, slide, title, categories, values,
                       x=0.6, y=1.5, w=12.0, h=5.5,
                       data_labels=True, number_format='0%'):
        """Add a doughnut chart (ring-shaped pie).

        Args:
            categories: list of segment labels.
            values: list of numeric values.
        """
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION

        chart_data = CategoryChartData()
        chart_data.categories = categories
        chart_data.add_series('Series 1', values)

        graphic_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.DOUGHNUT, Inches(x), Inches(y),
            Inches(w), Inches(h), chart_data)
        chart = graphic_frame.chart
        self._apply_chart_style(chart, has_legend=True)

        if data_labels:
            chart.plots[0].has_data_labels = True
            labels = chart.plots[0].data_labels
            labels.number_format = number_format
            labels.font.size = Pt(10)

        if title:
            chart.has_title = True
            chart.chart_title.text_frame.paragraphs[0].text = title
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)

        return chart
